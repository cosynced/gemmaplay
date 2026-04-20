"""Lesson content extraction.

Originally PDF-only (hence the module name, kept for backwards compatibility).
Now dispatches on suffix to the right extractor — text PDFs, scanned PDFs,
Word, PowerPoint, images of notes, plain text, and Markdown.

Scanned PDFs and photo uploads fall through to Gemma vision (via
``GemmaClient.extract_text_from_images``). That path is lazy-loaded so the
import graph doesn't drag in ``pdf2image`` / ``Pillow`` unless actually used.
"""
from __future__ import annotations

import asyncio
import hashlib
import tempfile
from pathlib import Path

import pdfplumber

from app.core.logging import get_logger

log = get_logger(__name__)


# Lowest number of characters we'll accept from pdfplumber before deciding a
# PDF is scanned and handing it to the vision pipeline. 100 chars comfortably
# clears a header-only first page but catches real image-only scans.
SCANNED_TEXT_THRESHOLD = 100

# Cap for rasterizing scanned PDFs — vision cost and latency grows with pages.
SCAN_PAGE_CAP = 10

SUPPORTED_SUFFIXES = {
    ".pdf", ".txt", ".md", ".docx", ".pptx", ".png", ".jpg", ".jpeg",
}


def extract_text(file_path: Path) -> str:
    """Backwards-compatible entry point.

    Only handles formats that don't require async (no vision calls). For the
    full dispatcher use :func:`extract_from_file`.
    """
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        # Legacy callers (existing orchestrator path) invoke this synchronously.
        # If a PDF turns out to be scanned, extract_from_file will handle it;
        # here we just return whatever text pdfplumber can get.
        return _extract_pdf_text_only(file_path)
    if suffix in {".txt", ".md"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")
    raise ValueError(f"Unsupported file type: {suffix}")


async def extract_from_file(file_path: Path) -> str:
    """Extract lesson text from any supported input format.

    Dispatches by suffix; falls back to Gemma vision for scanned PDFs and
    images. Returns a single text string ready to feed to the LessonAgent.
    """
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return await _extract_pdf(file_path)
    if suffix == ".docx":
        return _extract_docx(file_path)
    if suffix == ".pptx":
        return _extract_pptx(file_path)
    if suffix in {".png", ".jpg", ".jpeg"}:
        return await _extract_image(file_path)
    if suffix in {".txt", ".md"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")
    raise ValueError(f"Unsupported file type: {suffix}")


# ---------- PDF ----------

def _extract_pdf_text_only(file_path: Path) -> str:
    parts: list[str] = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            if text.strip():
                parts.append(text)
    return "\n\n".join(parts)


async def _extract_pdf(file_path: Path) -> str:
    parts: list[str] = []
    page_count = 0
    with pdfplumber.open(file_path) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            text = page.extract_text() or ""
            if text.strip():
                parts.append(text)
    text = "\n\n".join(parts)
    if len(text) >= SCANNED_TEXT_THRESHOLD:
        return text
    if page_count == 0:
        return text  # Nothing to rasterize.

    # Looks scanned — fall through to the vision pipeline.
    effective_pages = min(page_count, SCAN_PAGE_CAP)
    if page_count > SCAN_PAGE_CAP:
        log.warning(
            "scanned_pdf_page_cap",
            total_pages=page_count,
            using=effective_pages,
        )
    log.info(
        "scanned_pdf_detected",
        text_chars=len(text),
        pages=effective_pages,
    )
    return await _extract_via_vision_from_pdf(file_path, effective_pages)


async def _extract_via_vision_from_pdf(file_path: Path, page_limit: int) -> str:
    # Lazy imports: keep pdf2image/Pillow off the import graph when not used.
    from pdf2image import convert_from_path

    def _rasterize() -> list[Path]:
        tmp_dir = Path(tempfile.mkdtemp(prefix="gemmaplay_scan_"))
        images = convert_from_path(
            str(file_path), dpi=200, first_page=1, last_page=page_limit,
        )
        paths: list[Path] = []
        for i, img in enumerate(images):
            out = tmp_dir / f"page_{i + 1}.png"
            img.save(out, "PNG")
            paths.append(out)
        return paths

    image_paths = await asyncio.to_thread(_rasterize)
    try:
        from app.services.gemma_client import get_gemma_client

        return await get_gemma_client().extract_text_from_images(image_paths)
    finally:
        for p in image_paths:
            p.unlink(missing_ok=True)
        parent = image_paths[0].parent if image_paths else None
        if parent is not None:
            try:
                parent.rmdir()
            except OSError:
                pass


# ---------- Word ----------

def _extract_docx(file_path: Path) -> str:
    from docx import Document

    doc = Document(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n".join(paragraphs)


# ---------- PowerPoint ----------

def _extract_pptx(file_path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(file_path))
    slides: list[str] = []
    for slide in prs.slides:
        lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs)
                if text.strip():
                    lines.append(text)
        if lines:
            slides.append("\n".join(lines))
    return "\n\n".join(slides)


# ---------- Image ----------

async def _extract_image(file_path: Path) -> str:
    from app.services.gemma_client import get_gemma_client

    return await get_gemma_client().extract_text_from_images([file_path])


# ---------- Util ----------

def text_hash(text: str) -> str:
    """Stable hash for caching lesson extraction results."""
    return hashlib.sha256(text.encode("utf-8")).hexdigest()[:16]
