"""Cheap content-size inspection.

Lets the frontend warn users BEFORE we commit to the full extraction +
Gemma pipeline (which costs tokens and wall-clock). Never calls Gemma,
never persists anything.
"""
from __future__ import annotations

from dataclasses import asdict, dataclass
from pathlib import Path

import pdfplumber

MAX_SOURCE_CHARS = 12000
PDF_PAGE_CAP = 10
# Fallback char estimate per scanned page — about a page of handwritten notes.
SCANNED_PAGE_CHAR_ESTIMATE = 1500
SCANNED_TEXT_THRESHOLD = 100

SUPPORTED_SUFFIXES = {
    ".pdf", ".txt", ".md", ".docx", ".pptx", ".png", ".jpg", ".jpeg",
}


@dataclass
class InspectionResult:
    char_count: int
    will_truncate: bool
    truncate_at_chars: int
    estimated_pages: int | None
    pages_to_process: int
    file_type: str
    ok: bool

    def to_dict(self) -> dict:
        return asdict(self)


def inspect_text(text: str, *, file_type: str = "text") -> InspectionResult:
    n = len(text or "")
    return InspectionResult(
        char_count=n,
        will_truncate=n > MAX_SOURCE_CHARS,
        truncate_at_chars=MAX_SOURCE_CHARS,
        estimated_pages=None,
        pages_to_process=1,
        file_type=file_type,
        ok=True,
    )


def inspect_file(path: Path) -> InspectionResult:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        return InspectionResult(
            char_count=0,
            will_truncate=False,
            truncate_at_chars=MAX_SOURCE_CHARS,
            estimated_pages=None,
            pages_to_process=0,
            file_type=suffix.lstrip("."),
            ok=False,
        )
    if suffix == ".pdf":
        return _inspect_pdf(path)
    if suffix == ".docx":
        return _inspect_docx(path)
    if suffix == ".pptx":
        return _inspect_pptx(path)
    if suffix in {".png", ".jpg", ".jpeg"}:
        return InspectionResult(
            char_count=SCANNED_PAGE_CHAR_ESTIMATE,
            will_truncate=False,
            truncate_at_chars=MAX_SOURCE_CHARS,
            estimated_pages=1,
            pages_to_process=1,
            file_type=suffix.lstrip("."),
            ok=True,
        )
    # .txt / .md
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
    except OSError:
        text = ""
    r = inspect_text(text, file_type=suffix.lstrip("."))
    return r


def _inspect_pdf(path: Path) -> InspectionResult:
    try:
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)
            sampled_chars = 0
            sample_upper = min(total_pages, PDF_PAGE_CAP)
            for page in pdf.pages[:sample_upper]:
                try:
                    text = page.extract_text() or ""
                except Exception:  # noqa: BLE001
                    text = ""
                sampled_chars += len(text)
    except Exception:  # noqa: BLE001
        return InspectionResult(
            char_count=0,
            will_truncate=False,
            truncate_at_chars=MAX_SOURCE_CHARS,
            estimated_pages=None,
            pages_to_process=0,
            file_type="pdf",
            ok=False,
        )

    # If the sampled pages produced almost no text, it's a scanned PDF →
    # assume the vision path will extract roughly SCANNED_PAGE_CHAR_ESTIMATE
    # characters per page it processes.
    pages_to_process = min(total_pages, PDF_PAGE_CAP)
    if sampled_chars < SCANNED_TEXT_THRESHOLD and pages_to_process > 0:
        est_chars = pages_to_process * SCANNED_PAGE_CHAR_ESTIMATE
    else:
        est_chars = sampled_chars

    will_truncate = total_pages > PDF_PAGE_CAP or est_chars > MAX_SOURCE_CHARS

    return InspectionResult(
        char_count=est_chars,
        will_truncate=will_truncate,
        truncate_at_chars=MAX_SOURCE_CHARS,
        estimated_pages=total_pages,
        pages_to_process=pages_to_process,
        file_type="pdf",
        ok=True,
    )


def _inspect_docx(path: Path) -> InspectionResult:
    try:
        from docx import Document  # local import keeps startup lean
        doc = Document(str(path))
        chars = sum(len(p.text) for p in doc.paragraphs if p.text)
        paragraph_count = sum(1 for p in doc.paragraphs if p.text and p.text.strip())
    except Exception:  # noqa: BLE001
        return InspectionResult(
            char_count=0,
            will_truncate=False,
            truncate_at_chars=MAX_SOURCE_CHARS,
            estimated_pages=None,
            pages_to_process=0,
            file_type="docx",
            ok=False,
        )
    return InspectionResult(
        char_count=chars,
        will_truncate=chars > MAX_SOURCE_CHARS,
        truncate_at_chars=MAX_SOURCE_CHARS,
        estimated_pages=max(paragraph_count, 1),
        pages_to_process=max(paragraph_count, 1),
        file_type="docx",
        ok=True,
    )


def _inspect_pptx(path: Path) -> InspectionResult:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = list(prs.slides)
        slide_count = len(slides)
        chars = 0
        for slide in slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for para in shape.text_frame.paragraphs:
                    chars += sum(len(run.text or "") for run in para.runs)
    except Exception:  # noqa: BLE001
        return InspectionResult(
            char_count=0,
            will_truncate=False,
            truncate_at_chars=MAX_SOURCE_CHARS,
            estimated_pages=None,
            pages_to_process=0,
            file_type="pptx",
            ok=False,
        )
    return InspectionResult(
        char_count=chars,
        will_truncate=chars > MAX_SOURCE_CHARS or slide_count > PDF_PAGE_CAP,
        truncate_at_chars=MAX_SOURCE_CHARS,
        estimated_pages=slide_count,
        pages_to_process=min(slide_count, PDF_PAGE_CAP),
        file_type="pptx",
        ok=True,
    )
